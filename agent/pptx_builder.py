"""
agent/pptx_builder.py
--------------------
Renders a sales_deck deck_payload JSON spec into a .pptx bytes object.

Uses python-pptx. The Oracle color scheme:
  Red:   #C74634  (primary Oracle red)
  Dark:  #1A1A1A  (title backgrounds)
  Gray:  #F5F5F5  (content backgrounds)
  White: #FFFFFF

Called by: sub_agents/sales_deck/server.py after LLM returns deck_payload.
Returns: bytes — caller saves to object storage as .pptx.
"""
from __future__ import annotations

import io
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


ORACLE_RED = RGBColor(0xC7, 0x46, 0x34) if PPTX_AVAILABLE else None
ORACLE_DARK = RGBColor(0x1A, 0x1A, 0x1A) if PPTX_AVAILABLE else None
ORACLE_GRAY = RGBColor(0xF5, 0xF5, 0xF5) if PPTX_AVAILABLE else None
ORACLE_WHITE = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_AVAILABLE else None


def build_pptx(deck_payload: dict[str, Any]) -> bytes:
    """
    Render deck_payload into a .pptx and return the raw bytes.
    Raises ImportError if python-pptx is not installed.
    """
    if not PPTX_AVAILABLE:
        raise ImportError(
            "python-pptx is required for PPTX rendering. Run: pip install python-pptx"
        )

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for slide_spec in deck_payload.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        layout = slide_spec.get("layout", "content")
        title = slide_spec.get("title", "")
        notes = slide_spec.get("presenter_notes", "")

        _set_slide_background(slide, layout)
        _add_title_box(slide, title, layout)

        if layout == "title":
            subtitle = slide_spec.get("subtitle", "")
            if subtitle:
                _add_text_box(
                    slide,
                    subtitle,
                    Inches(1.5),
                    Inches(4.2),
                    Inches(10),
                    Inches(1),
                    font_size=24,
                    color=ORACLE_WHITE,
                    bold=False,
                )
        elif layout == "two_column":
            left = slide_spec.get("left_content", [])
            right = slide_spec.get("right_content", [])
            _add_bullet_box(slide, left, Inches(0.5), Inches(1.8), Inches(6), Inches(4.5))
            _add_bullet_box(slide, right, Inches(6.8), Inches(1.8), Inches(6), Inches(4.5))
        else:
            content = slide_spec.get("content", [])
            if isinstance(content, list):
                _add_bullet_box(
                    slide, content, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5)
                )
            elif isinstance(content, str) and content:
                _add_text_box(
                    slide, content, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5)
                )

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _set_slide_background(slide, layout: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = ORACLE_DARK if layout == "title" else ORACLE_WHITE


def _add_title_box(slide, text: str, layout: str) -> None:
    color = ORACLE_WHITE if layout == "title" else ORACLE_DARK
    tx_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.2)
    )
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32) if layout == "title" else Pt(28)
    p.font.bold = True
    p.font.color.rgb = color


def _add_text_box(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    font_size: int = 18,
    color=None,
    bold: bool = False,
) -> None:
    if color is None:
        color = ORACLE_DARK
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color


def _add_bullet_box(slide, items: list[str], left, top, width, height) -> None:
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = ORACLE_DARK
        p.space_after = Pt(6)
