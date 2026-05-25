"""Render a 7-slide Oracle OCI POC PowerPoint deck."""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from sub_agents.presentation.scripts.resolve_oci_powerpoint_icon import resolve_icon


TOOLKIT_PATH = (
    Path(__file__).parent.parent
    / "assets"
    / "oracle-oci-architecture-toolkit-v24.1.pptx"
)

ORACLE_RED = RGBColor(196, 0, 0)
ORACLE_DARK = RGBColor(49, 45, 42)
ORACLE_GRAY = RGBColor(245, 244, 242)
WHITE = RGBColor(255, 255, 255)


def render(spec: dict[str, Any], output_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, spec)
    _add_challenge_slide(prs, spec)
    _add_architecture_slide(prs, spec)
    _add_services_slide(prs, spec)
    _add_bom_slide(prs, spec)
    _add_timeline_slide(prs, spec)
    _add_next_steps_slide(prs, spec)

    prs.save(output_path)


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def _add_header(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ORACLE_DARK
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(0.95),
        Inches(12.2),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ORACLE_RED
    line.line.fill.background()


def _text(slide, text: str, left, top, width, height, size: int = 18, bold: bool = False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text or "")
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = ORACLE_DARK
    return box


def _bullet_list(slide, items: list[Any], left, top, width, height, size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    safe_items = [str(item) for item in items if str(item).strip()]
    for idx, item in enumerate(safe_items or ["Details to be confirmed."]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = ORACLE_DARK


def _add_title_slide(prs: Presentation, spec: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(1.15),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ORACLE_RED
    band.line.fill.background()
    _text(slide, "Oracle Cloud Infrastructure POC", Inches(0.65), Inches(0.32), Inches(8), Inches(0.5), 24, True).text_frame.paragraphs[0].font.color.rgb = WHITE
    _text(slide, str(spec.get("poc_name") or "OCI POC"), Inches(0.85), Inches(2.0), Inches(11.6), Inches(0.9), 36, True)
    _text(slide, str(spec.get("customer_name") or "Customer"), Inches(0.85), Inches(3.05), Inches(9.5), Inches(0.45), 22)
    _text(slide, str(spec.get("date") or ""), Inches(0.85), Inches(3.65), Inches(5), Inches(0.35), 16)


def _add_challenge_slide(prs: Presentation, spec: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "Customer Challenge")
    pain = str(spec.get("pain_statement") or "Business pain and success criteria to be confirmed.")
    _text(slide, pain, Inches(0.85), Inches(1.55), Inches(11.7), Inches(1.0), 24, True)
    facts = spec.get("customer_facts") or spec.get("success_criteria") or []
    if not isinstance(facts, list):
        facts = [facts]
    _bullet_list(slide, facts, Inches(1.0), Inches(3.0), Inches(11), Inches(2.7), 18)


def _add_architecture_slide(prs: Presentation, spec: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "OCI Architecture")
    services = _service_names(spec)
    if not services:
        services = ["Virtual Cloud Network", "OCI Compute", "Autonomous Database"]
    max_services = min(len(services), 5)
    spacing = 11.0 / max(max_services, 1)
    for idx, service in enumerate(services[:5]):
        left = Inches(0.95 + idx * spacing)
        top = Inches(2.35)
        shape_name = resolve_icon(service)
        copied = False
        if shape_name:
            copied = _copy_icon_from_toolkit(
                str(TOOLKIT_PATH),
                shape_name,
                slide,
                left,
                top,
                Inches(1.05),
                Inches(1.05),
            )
        if not copied:
            _add_fallback_icon(slide, service, left, top)
        label = _text(slide, service, left - Inches(0.18), Inches(3.55), Inches(1.45), Inches(0.6), 11, True)
        label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _text(slide, "Architecture uses official OCI service names and copies native toolkit shapes when available.", Inches(0.95), Inches(5.3), Inches(11.5), Inches(0.45), 15)


def _add_services_slide(prs: Presentation, spec: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "Key OCI Services")
    services = _service_names(spec)
    _bullet_list(slide, services, Inches(1.0), Inches(1.45), Inches(11.3), Inches(5.3), 18)


def _add_bom_slide(prs: Presentation, spec: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "Cost Estimate")
    bom_summary = str(spec.get("bom_summary") or "BOM estimate to be finalized after sizing confirmation.")
    _text(slide, bom_summary, Inches(0.85), Inches(1.4), Inches(11.6), Inches(0.85), 20, True)
    rows = spec.get("bom_rows") or []
    if not isinstance(rows, list):
        rows = []
    items = []
    for row in rows[:6]:
        if isinstance(row, dict):
            items.append(f"{row.get('service', 'OCI service')}: ${row.get('monthly_cost', 'TBD')}/mo")
    _bullet_list(slide, items or ["Monthly estimate: TBD"], Inches(1.0), Inches(2.7), Inches(11), Inches(3.5), 17)


def _add_timeline_slide(prs: Presentation, spec: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "Implementation Timeline")
    phases = spec.get("jep_phases") or ["Prepare", "Build", "Validate", "Executive readout"]
    if not isinstance(phases, list):
        phases = [phases]
    width = 11.0 / max(len(phases[:5]), 1)
    for idx, phase in enumerate(phases[:5]):
        left = Inches(0.85 + idx * width)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            Inches(2.25),
            Inches(width - 0.18),
            Inches(1.25),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = ORACLE_GRAY
        card.line.color.rgb = ORACLE_RED
        p = card.text_frame.paragraphs[0]
        p.text = f"{idx + 1}. {phase}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ORACLE_DARK


def _add_next_steps_slide(prs: Presentation, spec: dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "Next Steps")
    steps = spec.get("next_steps") or [
        "Confirm POC option and success criteria",
        "Validate region, tenancy access, and prerequisite quotas",
        "Schedule build window and customer demo",
        "Run executive readout with measured outcome",
    ]
    if not isinstance(steps, list):
        steps = [steps]
    _bullet_list(slide, steps, Inches(1.0), Inches(1.55), Inches(11.3), Inches(4.8), 20)


def _copy_icon_from_toolkit(
    toolkit_pptx_path: str,
    shape_name: str,
    target_slide,
    left,
    top,
    width,
    height,
) -> bool:
    path = Path(toolkit_pptx_path)
    if not path.exists():
        return False
    toolkit = Presentation(str(path))
    for slide in toolkit.slides:
        for shape in slide.shapes:
            if shape.name == shape_name:
                sp = deepcopy(shape._element)
                target_slide.shapes._spTree.append(sp)
                shape_obj = target_slide.shapes[-1]
                shape_obj.left = left
                shape_obj.top = top
                shape_obj.width = width
                shape_obj.height = height
                return True
    return False


def _add_fallback_icon(slide, service: str, left, top) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        Inches(1.05),
        Inches(1.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORACLE_GRAY
    shape.line.color.rgb = ORACLE_RED
    p = shape.text_frame.paragraphs[0]
    p.text = _initials(service)
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = ORACLE_DARK
    p.alignment = PP_ALIGN.CENTER


def _service_names(spec: dict[str, Any]) -> list[str]:
    services = spec.get("oci_services") or []
    if not isinstance(services, list):
        services = [services]
    names: list[str] = []
    for service in services:
        if isinstance(service, dict):
            value = service.get("name") or service.get("service") or service.get("display_name")
        else:
            value = service
        text = str(value or "").strip()
        if text:
            names.append(text)
    return names


def _initials(value: str) -> str:
    words = [word for word in str(value or "").replace("-", " ").split() if word]
    return "".join(word[0].upper() for word in words[:3]) or "OCI"
